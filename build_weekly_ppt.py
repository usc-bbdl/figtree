import os.path
from datetime import datetime, timedelta
from os import mkdir, listdir
from shutil import copy2

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Emu

### Assumes files will be downloaded to a folder name TEMP_INPUT_FOLDER
TEMP_INPUT_FOLDER = 'TEMP_INPUT_FOLDER/'


def build_weekly_ppt(TEMP_INPUT_FOLDER):
    """
    Takes in the name of the folder that contains the 'Figure Queue' as well as the 'Meeting Agenda' (TEMP_INPUT_FOLDER), builds a powerpoint presentation for that week's lab meeting, and moves the presentation and corresponding figures to a new folder that will be uploaded to Google Drive (TEMP_OUTPUT_FOLDER)
    """

    if TEMP_INPUT_FOLDER[-1] != "/": TEMP_INPUT_FOLDER = TEMP_INPUT_FOLDER + "/"
    assert os.path.exists(TEMP_INPUT_FOLDER)

    # TEMP_OUTPUT_FOLDER = "TEMP_OUTPUT_FOLDER"
    # mkdir(TEMP_OUTPUT_FOLDER)

    todaysDate = datetime.now()
    daysToLabMeeting = 0
    date = todaysDate
    for i in range(7):
        if date.weekday() == 1:  # Tuesday
            daysToLabMeeting = i
            break
        else:
            date = datetime.now() + timedelta(days=i + 1)
    labMeetingDate = datetime.now() + timedelta(days=daysToLabMeeting)
    labMeetingDateStr = labMeetingDate.strftime("%A, %B %d, %Y")
    labMeetingFolderName = labMeetingDate.strftime("%Y_%m_%d Lab Meeting/")
    mkdir(labMeetingFolderName)
    mkdir(labMeetingFolderName + "Figures/")
    labMeetingPresentationName = labMeetingDate.strftime(
        "%Y_%m_%d_Lab_Meeting.pptx"
    )

    # Find all figures in Figure Queue
    fileNames = [
        f
        for f in listdir(TEMP_INPUT_FOLDER + "Figure Queue/")
        if os.path.isfile(os.path.join(TEMP_INPUT_FOLDER + "Figure Queue/", f))
    ]

    # Create presentation from BBDL format
    prs = Presentation('graphics_and_templates/template.pptx')
    title_slide_layout = prs.slide_layouts[0]
    agenda_slide_layout = prs.slide_layouts[1]
    content_slide_layout = prs.slide_layouts[2]

    # Add Title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide_title = title_slide.placeholders[10]
    title_slide_subtitle = title_slide.placeholders[11]
    title_slide_title.text = "Valero Lab Weekly Meeting"
    title_slide_subtitle.text = labMeetingDateStr

    # Add Agenda slide
    agenda_xlsx = pd.read_excel(
        TEMP_INPUT_FOLDER+"Weekly Agenda.xlsx",
        sheet_name="Formatted Agenda Items"
    )
    agendaItems=list(agenda_xlsx["Agenda Item(s)"])[0]
    whitelist = {
        "and", "as", "at", "but", "by", "for", "from", "if", "in", "into",
        "like", "near", "nor",  "of", "off", "on", "once", "onto", "or",
        "over", "past", "so", "than", "that", "till", "to", "up", "upon",
        "with", "when", "yet", "is", "a", "an", "the"
    }
    if type(agendaItems)==float:
        agendaItemsCount = 0
    if type(agendaItems)==str:
        agendaItems = agendaItems.split(";")
        for i in range(len(agendaItems)):
            # remove leading/trailing whitespace and capitalize all words
            temp=agendaItems[i].lstrip().rstrip().lower().split(" ")
            agendaItems[i]=temp[0].capitalize()
            for j in range(1,len(temp)):
                if temp[j] not in whitelist:
                    agendaItems[i]+= " " + temp[j].capitalize()
                else:
                    agendaItems[i]+= " " + temp[j]
        agendaItemsCount = len(agendaItems)

    agenda_slide = prs.slides.add_slide(agenda_slide_layout)
    agenda_slide_title = agenda_slide.placeholders[10]
    agenda_slide_subtitle = agenda_slide.placeholders[11]
    agenda_slide_title.text = "Today's Meeting Agenda"
    if agendaItemsCount==0:
        thisWeeksAgenda = "No Agenda Items Were Added This Week."
    else:
        thisWeeksAgenda = ""
        for item in agendaItems:
            thisWeeksAgenda += " - " + item + "\n"
    agenda_slide_subtitle.text = thisWeeksAgenda

    # Add all appropriate files as Content slides
    figureCount = 0
    for file in fileNames:
        content_slide = prs.slides.add_slide(content_slide_layout)
        content_slide_title = content_slide.placeholders[10]
        if "_" not in str(file):
            titleStr = str(file)
            presenterName = "Unknown"
        else:
            titleStr = file.replace("_", " ")[:file.find(".")]
            presenterName = titleStr[(titleStr.rfind(" ") + 1):]
            titleStr = titleStr[:titleStr.rfind(" ")]
        content_slide_title.text = titleStr + " (" + presenterName + ")"
        figure = content_slide.shapes.add_picture(
            TEMP_INPUT_FOLDER + "Figure Queue/" + file,
            Inches(0.5),
            Inches(1.75)
        )
        ratio = figure.height / figure.width
        if ratio > 5 / 13:
            figure.height = Emu(Inches(5))
            figure.width = Emu(Inches(5 / ratio))
        else:
            figure.width = Emu(Inches(13))
            figure.height = Emu(Inches(13 * ratio))
        figure.left = Emu((Inches(13.3333) - figure.width) / 2)
        figure.top = Emu((Inches(7.5) - figure.height) / 2)
        copy2(
            TEMP_INPUT_FOLDER + "Figure Queue/" + file,
            labMeetingFolderName + "Figures/"
        )
        figureCount += 1

    # Remove template slides
    rIdToDrop = ['rId1', 'rId2', 'rId3', 'rId4']
    for i in reversed(range(len(prs.slides._sldIdLst))):
        if prs.slides._sldIdLst[i].rId in rIdToDrop:
            prs.part.drop_rel(prs.slides._sldIdLst[i].rId)
            del prs.slides._sldIdLst[i]

    prs.save(labMeetingPresentationName)
    return (agendaItemsCount, figureCount, labMeetingFolderName)


assert os.path.exists(TEMP_INPUT_FOLDER), "Need to create TEMP_INPUT_FOLDER."

agendaItemsCount, figureCount, labMeetingFolderName = build_weekly_ppt(
    TEMP_INPUT_FOLDER
)

print("agendaItemsCount\t" + str(agendaItemsCount) + "\n"
    + "figureCount\t" + str(figureCount) + "\n"
    + "labMeetingFolderName\t" + labMeetingFolderName
)
